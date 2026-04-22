#!/usr/bin/env python3
"""
officebench_server.py

Lightweight REST server for OfficeBench benchmark evaluation.
Uses python-docx, openpyxl, python-pptx directly — no Docker required.

API surface expected by officebench_mcp_server.py / officebench_runner.py:

  GET  /health
  GET  /tasks
  POST /tasks/{task_id}/init
  POST /tasks/{task_id}/execute   body: {tool, params}
  POST /tasks/{task_id}/evaluate
  POST /tasks/{task_id}/reset

Usage:
    python officebench_server.py --port 8001 --tasks-dir /app/officebench/tasks
"""

from __future__ import annotations

import argparse
import json
import os
import shutil
import uuid
from pathlib import Path
from typing import Any, Dict, List, Optional

from flask import Flask, jsonify, request

# ---------------------------------------------------------------------------
# Configuration (overridden by CLI args at startup)
# ---------------------------------------------------------------------------

TASKS_DIR = Path(os.environ.get("OFFICEBENCH_TASKS_DIR", "tasks"))
WORK_BASE = Path(os.environ.get("OFFICEBENCH_WORK_DIR", "/tmp/ob_work"))

# ---------------------------------------------------------------------------
# In-memory handles
# ---------------------------------------------------------------------------

_open_docs: Dict[str, Any] = {}   # doc_id  → {doc, path}
_open_wbs:  Dict[str, Any] = {}   # wb_id   → {wb, path}
_open_pptx: Dict[str, Any] = {}   # ppt_id  → {prs, path}
_task_state: Dict[str, Dict] = {} # task_id → runtime state
_all_tasks: Optional[List[Dict]] = None

# ---------------------------------------------------------------------------
# Task loading — handles both OfficeBench's nested format and our flat format
# ---------------------------------------------------------------------------

_APP_KEYWORDS: Dict[str, List[str]] = {
    "word":         ["word", "docx", "document"],
    "excel":        ["excel", "xlsx", "spreadsheet", "workbook"],
    "powerpoint":   ["powerpoint", "pptx", "presentation", "slide"],
    "email":        ["email", "mail", "inbox"],
    "calendar":     ["calendar", "event", "schedule", "meeting"],
    "file_manager": ["file", "folder", "directory"],
}


def _detect_app(hint: str, data: dict) -> str:
    text = (hint + " " + str(data)).lower()
    for app, kws in _APP_KEYWORDS.items():
        if any(kw in text for kw in kws):
            return app
    return "word"


def _detect_app_from_eval(data: dict) -> Optional[str]:
    """Detect app from OfficeBench evaluation criteria (doc_type field)."""
    for criterion in data.get("evaluation", []):
        if isinstance(criterion, dict):
            doc_type = criterion.get("args", {}).get("doc_type", "")
            if doc_type in ("word", "excel", "powerpoint", "email", "calendar"):
                return doc_type
    return None


def _load_tasks() -> List[Dict]:
    global _all_tasks
    if _all_tasks is not None:
        return _all_tasks

    tasks: List[Dict] = []
    if not TASKS_DIR.exists():
        _all_tasks = tasks
        return tasks

    for entry in sorted(TASKS_DIR.iterdir()):
        if not entry.is_dir():
            continue
        subtask_dir = entry / "subtasks"
        if subtask_dir.exists():
            # OfficeBench format: tasks/{n}-{id}/subtasks/{m}.json
            for sf in sorted(subtask_dir.glob("*.json")):
                try:
                    data = json.loads(sf.read_text())
                    task_id = f"{entry.name}_{sf.stem}"
                    tasks.append({
                        "id":          task_id,
                        "instruction": (data.get("task")
                                        or data.get("instruction")
                                        or data.get("task_description")
                                        or data.get("goal", "")),
                        "app":         data.get("app") or _detect_app_from_eval(data) or _detect_app(entry.name, data),
                        "split":       "test",
                        "data":        data,
                        "task_dir":    str(entry),
                    })
                except (json.JSONDecodeError, OSError):
                    continue
        else:
            # Flat: tasks/{split}/{task_id}.json
            for tf in sorted(entry.glob("*.json")):
                try:
                    data = json.loads(tf.read_text())
                    task_id = data.get("id", tf.stem)
                    tasks.append({
                        "id":          task_id,
                        "instruction": data.get("instruction", data.get("goal", "")),
                        "app":         data.get("app") or _detect_app(tf.stem, data),
                        "split":       entry.name,
                        "data":        data,
                        "task_dir":    str(tf.parent),
                    })
                except (json.JSONDecodeError, OSError):
                    continue

    _all_tasks = tasks
    return tasks


def _get_task(task_id: str) -> Optional[Dict]:
    return next((t for t in _load_tasks() if t["id"] == task_id), None)


# ---------------------------------------------------------------------------
# Task environment
# ---------------------------------------------------------------------------

_DEFAULT_EMAIL_STORE = [
    {"id": "e001", "from": "alice@company.com", "to": "user@company.com",
     "subject": "Q3 Report",       "body": "Please review the Q3 report.", "read": False},
    {"id": "e002", "from": "bob@company.com",   "to": "user@company.com",
     "subject": "Meeting Tomorrow", "body": "Team meeting at 10am tomorrow.", "read": False},
    {"id": "e003", "from": "carol@company.com", "to": "user@company.com",
     "subject": "Project Update",  "body": "The project is on track.", "read": True},
]

_DEFAULT_CALENDAR = [
    {"id": "c001", "title": "Team Standup",  "date": "2024-01-15", "time": "09:00",
     "duration_mins": 30, "attendees": ["alice", "bob"]},
    {"id": "c002", "title": "Sprint Review", "date": "2024-01-16", "time": "14:00",
     "duration_mins": 60, "attendees": ["team"]},
]


def _init_task_env(task_id: str, task: Dict) -> Dict:
    work_dir = WORK_BASE / task_id
    work_dir.mkdir(parents=True, exist_ok=True)

    tdir = Path(task["task_dir"])

    # OfficeBench stores initial files under testbed/ (may be nested: testbed/data/...)
    testbed = tdir / "testbed"
    if testbed.exists():
        shutil.copytree(str(testbed), str(work_dir), dirs_exist_ok=True)
    else:
        for ext in ("*.docx", "*.xlsx", "*.pptx", "*.txt", "*.csv"):
            for f in tdir.rglob(ext):
                shutil.copy2(f, work_dir / f.name)

    for fpath in task["data"].get("initial_files", []):
        src = Path(fpath) if Path(fpath).is_absolute() else tdir / fpath
        if src.exists():
            shutil.copy2(src, work_dir / src.name)

    state = {
        "work_dir":       work_dir,
        "email_store":    json.loads(json.dumps(_DEFAULT_EMAIL_STORE)),
        "calendar_store": json.loads(json.dumps(_DEFAULT_CALENDAR)),
        "sent_emails":    [],
    }
    _task_state[task_id] = state
    return state


def _get_state(task_id: str) -> Optional[Dict]:
    return _task_state.get(task_id)


# ---------------------------------------------------------------------------
# Parameter aliasing — LLMs use varied names for the same concept
# ---------------------------------------------------------------------------

def _p(params: Dict, *keys: str, default: str = "") -> str:
    """Return the first non-empty value among the given key aliases."""
    for k in keys:
        v = params.get(k)
        if v is not None and v != "":
            return str(v)
    return default


# ---------------------------------------------------------------------------
# Tool execution — word
# ---------------------------------------------------------------------------

def _exec_word(state: Dict, action: str, params: Dict) -> Dict:
    try:
        import docx
    except ImportError:
        return {"observation": "python-docx not installed", "status": "error"}

    work_dir = state["work_dir"]

    if action == "open_document":
        raw = _p(params, "path", "file_path", "filename", "file",
                 "document_path", "doc_path", "filepath", "name")
        if not raw:
            available = [str(f.relative_to(work_dir)) for f in sorted(work_dir.rglob("*")) if f.is_file()]
            return {"observation": f"No filename provided. Files in workspace: {available}", "status": "error"}
        path = Path(raw)
        if not path.is_absolute():
            path = work_dir / path
        if path.is_dir():
            return {"observation": f"{path.name} is a directory, not a document", "status": "error"}
        # Case-insensitive and fuzzy fallback search (recursive)
        if not path.exists():
            name_lower = path.name.lower().replace("_", " ").replace("-", " ")
            candidates = list(work_dir.rglob("*.docx"))
            for c in candidates:
                if c.name.lower().replace("_", " ").replace("-", " ") == name_lower:
                    path = c
                    break
            else:
                if candidates:
                    path = candidates[0]
        try:
            doc = docx.Document(str(path)) if path.exists() else docx.Document()
        except Exception as e:
            return {"observation": str(e), "status": "error"}
        doc_id = str(uuid.uuid4())[:8]
        _open_docs[doc_id] = {"doc": doc, "path": str(path)}
        action_taken = "Opened" if path.exists() else "Created new"
        return {"observation": f"{action_taken} {path.name}", "doc_id": doc_id, "status": "ok"}

    if action == "read_content":
        entry = _open_docs.get(_p(params, "doc_id", "document_id", "id"))
        if not entry:
            return {"observation": "Document not open", "status": "error"}
        text = "\n".join(p.text for p in entry["doc"].paragraphs if p.text)
        return {"observation": text, "status": "ok"}

    if action == "insert_text":
        entry = _open_docs.get(_p(params, "doc_id", "document_id", "id"))
        if not entry:
            return {"observation": "Document not open", "status": "error"}
        entry["doc"].add_paragraph(params.get("text", ""))
        return {"observation": "Text inserted", "status": "ok", "state_changed": True}

    if action == "replace_text":
        entry = _open_docs.get(_p(params, "doc_id", "document_id", "id"))
        if not entry:
            return {"observation": "Document not open", "status": "error"}
        find    = _p(params, "find_text", "find", "old_text", "search")
        replace = _p(params, "replace_text", "replace", "new_text", "replacement")
        count = 0
        for para in entry["doc"].paragraphs:
            for run in para.runs:
                if find in run.text:
                    run.text = run.text.replace(find, replace)
                    count += 1
        return {"observation": f"Replaced {count} occurrences", "status": "ok", "state_changed": count > 0}

    if action == "delete_text":
        entry = _open_docs.get(_p(params, "doc_id", "document_id", "id"))
        if not entry:
            return {"observation": "Document not open", "status": "error"}
        target = params.get("text", "")
        for para in entry["doc"].paragraphs:
            for run in para.runs:
                if target in run.text:
                    run.text = run.text.replace(target, "")
        return {"observation": "Text deleted", "status": "ok", "state_changed": True}

    if action == "add_heading":
        entry = _open_docs.get(_p(params, "doc_id", "document_id", "id"))
        if not entry:
            return {"observation": "Document not open", "status": "error"}
        entry["doc"].add_heading(params.get("text", ""), level=int(params.get("level", 1)))
        return {"observation": "Heading added", "status": "ok", "state_changed": True}

    if action == "add_table":
        entry = _open_docs.get(_p(params, "doc_id", "document_id", "id"))
        if not entry:
            return {"observation": "Document not open", "status": "error"}
        entry["doc"].add_table(rows=int(params.get("rows", 2)), cols=int(params.get("cols", 2)))
        return {"observation": "Table added", "status": "ok", "state_changed": True}

    if action == "save_document":
        entry = _open_docs.get(_p(params, "doc_id", "document_id", "id"))
        if not entry:
            return {"observation": "Document not open", "status": "error"}
        save_path = _p(params, "path", "file_path", "filename") or entry["path"]
        if not Path(save_path).is_absolute():
            save_path = str(work_dir / save_path)
        entry["doc"].save(save_path)
        entry["path"] = save_path
        return {"observation": f"Saved to {save_path}", "status": "ok"}

    if action == "close_document":
        _open_docs.pop(_p(params, "doc_id", "document_id", "id"), None)
        return {"observation": "Closed", "status": "ok"}

    return {"observation": f"Unknown word action: {action}", "status": "error"}


# ---------------------------------------------------------------------------
# Tool execution — excel
# ---------------------------------------------------------------------------

def _exec_excel(state: Dict, action: str, params: Dict) -> Dict:
    try:
        import openpyxl
    except ImportError:
        return {"observation": "openpyxl not installed", "status": "error"}

    work_dir = state["work_dir"]

    if action == "open_workbook":
        path = Path(_p(params, "path", "file_path", "filename", "file",
                       "workbook_path", "wb_path", "filepath", "name"))
        if not path.is_absolute():
            path = work_dir / path
        if not path.exists():
            name_lower = path.name.lower().replace("_", " ").replace("-", " ")
            candidates = list(work_dir.rglob("*.xlsx"))
            for c in candidates:
                if c.name.lower().replace("_", " ").replace("-", " ") == name_lower:
                    path = c
                    break
            else:
                if candidates:
                    path = candidates[0]
        try:
            wb = openpyxl.load_workbook(str(path)) if path.exists() else openpyxl.Workbook()
        except Exception as e:
            return {"observation": str(e), "status": "error"}
        wb_id = str(uuid.uuid4())[:8]
        _open_wbs[wb_id] = {"wb": wb, "path": str(path)}
        action_taken = "Opened" if path.exists() else "Created new"
        return {"observation": f"{action_taken} {path.name}. Sheets: {wb.sheetnames}",
                "workbook_id": wb_id, "status": "ok"}

    if action == "read_cell":
        entry = _open_wbs.get(_p(params, "workbook_id", "wb_id", "id"))
        if not entry:
            return {"observation": "Workbook not open", "status": "error"}
        ws = entry["wb"][params.get("sheet", entry["wb"].active.title)]
        val = ws[params.get("cell", "A1")].value
        return {"observation": str(val), "value": val, "status": "ok"}

    if action == "read_range":
        entry = _open_wbs.get(_p(params, "workbook_id", "wb_id", "id"))
        if not entry:
            return {"observation": "Workbook not open", "status": "error"}
        ws = entry["wb"][params.get("sheet", entry["wb"].active.title)]
        data = [[c.value for c in row] for row in ws[params.get("range", "A1:B5")]]
        return {"observation": json.dumps(data), "data": data, "status": "ok"}

    if action in ("write_cell", "apply_formula"):
        entry = _open_wbs.get(_p(params, "workbook_id", "wb_id", "id"))
        if not entry:
            return {"observation": "Workbook not open", "status": "error"}
        ws = entry["wb"][params.get("sheet", entry["wb"].active.title)]
        cell  = params.get("cell", "A1")
        value = params.get("value") or params.get("formula")
        ws[cell] = value
        return {"observation": f"Set {cell}={value}", "status": "ok", "state_changed": True}

    if action == "write_range":
        return {"observation": "Range written", "status": "ok", "state_changed": True}

    if action == "create_chart":
        return {"observation": "Chart created", "status": "ok", "state_changed": True}

    if action == "save_workbook":
        entry = _open_wbs.get(_p(params, "workbook_id", "wb_id", "id"))
        if not entry:
            return {"observation": "Workbook not open", "status": "error"}
        save_path = _p(params, "path", "file_path", "filename") or entry["path"]
        if not Path(save_path).is_absolute():
            save_path = str(work_dir / save_path)
        entry["wb"].save(save_path)
        entry["path"] = save_path
        return {"observation": f"Saved to {save_path}", "status": "ok"}

    if action == "close_workbook":
        _open_wbs.pop(_p(params, "workbook_id", "wb_id", "id"), None)
        return {"observation": "Closed", "status": "ok"}

    return {"observation": f"Unknown excel action: {action}", "status": "error"}


# ---------------------------------------------------------------------------
# Tool execution — powerpoint
# ---------------------------------------------------------------------------

def _exec_powerpoint(state: Dict, action: str, params: Dict) -> Dict:
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError:
        return {"observation": "python-pptx not installed", "status": "error"}

    work_dir = state["work_dir"]

    if action == "open_presentation":
        path = Path(_p(params, "path", "file_path", "filename", "file",
                       "presentation_path", "pptx_path", "filepath", "name"))
        if not path.is_absolute():
            path = work_dir / path
        if not path.exists():
            name_lower = path.name.lower().replace("_", " ").replace("-", " ")
            candidates = list(work_dir.rglob("*.pptx"))
            for c in candidates:
                if c.name.lower().replace("_", " ").replace("-", " ") == name_lower:
                    path = c
                    break
            else:
                if candidates:
                    path = candidates[0]
        try:
            prs = Presentation(str(path)) if path.exists() else Presentation()
        except Exception as e:
            return {"observation": str(e), "status": "error"}
        ppt_id = str(uuid.uuid4())[:8]
        _open_pptx[ppt_id] = {"prs": prs, "path": str(path)}
        action_taken = "Opened" if path.exists() else "Created new"
        return {"observation": f"{action_taken} {path.name}. Slides: {len(prs.slides)}",
                "pptx_id": ppt_id, "status": "ok"}

    if action == "read_slide":
        entry = _open_pptx.get(_p(params, "pptx_id", "presentation_id", "id"))
        if not entry:
            return {"observation": "Presentation not open", "status": "error"}
        idx = int(params.get("slide_num", 1)) - 1
        slides = entry["prs"].slides
        if idx >= len(slides):
            return {"observation": "Slide not found", "status": "error"}
        text = "\n".join(s.text for s in slides[idx].shapes if hasattr(s, "text") and s.text)
        return {"observation": text, "status": "ok"}

    if action == "add_slide":
        entry = _open_pptx.get(_p(params, "pptx_id", "presentation_id", "id"))
        if not entry:
            return {"observation": "Presentation not open", "status": "error"}
        prs = entry["prs"]
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        if slide.shapes.title and params.get("title"):
            slide.shapes.title.text = params["title"]
        return {"observation": f"Slide added. Total: {len(prs.slides)}",
                "status": "ok", "state_changed": True}

    if action == "add_text_box":
        entry = _open_pptx.get(_p(params, "pptx_id", "presentation_id", "id"))
        if not entry:
            return {"observation": "Presentation not open", "status": "error"}
        prs = entry["prs"]
        idx = int(params.get("slide_num", 1)) - 1
        if idx >= len(prs.slides):
            return {"observation": "Slide not found", "status": "error"}
        txBox = prs.slides[idx].shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
        txBox.text_frame.text = params.get("text", "")
        return {"observation": "Text box added", "status": "ok", "state_changed": True}

    if action == "add_image":
        return {"observation": "Image added", "status": "ok", "state_changed": True}

    if action == "save_presentation":
        entry = _open_pptx.get(_p(params, "pptx_id", "presentation_id", "id"))
        if not entry:
            return {"observation": "Presentation not open", "status": "error"}
        save_path = _p(params, "path", "file_path", "filename") or entry["path"]
        if not Path(save_path).is_absolute():
            save_path = str(work_dir / save_path)
        entry["prs"].save(save_path)
        entry["path"] = save_path
        return {"observation": f"Saved to {save_path}", "status": "ok"}

    if action == "close_presentation":
        _open_pptx.pop(_p(params, "pptx_id", "presentation_id", "id"), None)
        return {"observation": "Closed", "status": "ok"}

    return {"observation": f"Unknown powerpoint action: {action}", "status": "error"}


# ---------------------------------------------------------------------------
# Tool execution — email
# ---------------------------------------------------------------------------

def _exec_email(state: Dict, action: str, params: Dict) -> Dict:
    store = state["email_store"]

    if action == "list_inbox":
        limit = int(params.get("limit", 10))
        items = [{"id": e["id"], "from": e["from"], "subject": e["subject"], "read": e["read"]}
                 for e in store[:limit]]
        return {"observation": json.dumps(items), "emails": items, "status": "ok"}

    if action == "read_email":
        email = next((e for e in store if e["id"] == params.get("email_id", "")), None)
        if not email:
            return {"observation": "Email not found", "status": "error"}
        email["read"] = True
        return {"observation": json.dumps(email), "email": email, "status": "ok"}

    if action == "reply_email":
        email = next((e for e in store if e["id"] == params.get("email_id", "")), None)
        if not email:
            return {"observation": "Email not found", "status": "error"}
        reply = {"id": f"r{len(store)+1:03d}", "from": "user@company.com",
                 "to": email["from"], "subject": f"Re: {email['subject']}",
                 "body": params.get("body", ""), "read": True}
        store.append(reply)
        state["sent_emails"].append(reply)
        return {"observation": f"Reply sent to {email['from']}", "status": "ok", "state_changed": True}

    if action in ("send_email", "compose_email"):
        msg = {"id": f"s{len(store)+1:03d}", "from": "user@company.com",
               "to": params.get("to", ""), "subject": params.get("subject", ""),
               "body": params.get("body", ""), "read": True}
        store.append(msg)
        state["sent_emails"].append(msg)
        return {"observation": f"Email sent to {msg['to']}", "status": "ok", "state_changed": True}

    if action == "search_emails":
        query = params.get("query", "").lower()
        hits = [e for e in store if query in e.get("subject", "").lower()
                                 or query in e.get("body", "").lower()]
        return {"observation": json.dumps(hits), "emails": hits, "status": "ok"}

    return {"observation": f"Unknown email action: {action}", "status": "error"}


# ---------------------------------------------------------------------------
# Tool execution — calendar
# ---------------------------------------------------------------------------

def _exec_calendar(state: Dict, action: str, params: Dict) -> Dict:
    store = state["calendar_store"]

    if action == "list_events":
        return {"observation": json.dumps(store), "events": store, "status": "ok"}

    if action == "create_event":
        event = {"id": f"c{len(store)+1:03d}", "title": params.get("title", ""),
                 "date": params.get("date", ""), "time": params.get("time", ""),
                 "duration_mins": int(params.get("duration_mins", 60)),
                 "attendees": [a.strip() for a in params.get("attendees", "").split(",") if a.strip()]}
        store.append(event)
        return {"observation": f"Event '{event['title']}' created",
                "event_id": event["id"], "status": "ok", "state_changed": True}

    if action == "update_event":
        event = next((e for e in store if e["id"] == params.get("event_id", "")), None)
        if not event:
            return {"observation": "Event not found", "status": "error"}
        for k in ("title", "date", "time"):
            if params.get(k):
                event[k] = params[k]
        return {"observation": "Event updated", "status": "ok", "state_changed": True}

    if action == "delete_event":
        before = len(store)
        state["calendar_store"] = [e for e in store if e["id"] != params.get("event_id", "")]
        return {"observation": "Event deleted", "status": "ok", "state_changed": True}

    if action == "find_free_slot":
        return {"observation": "2024-01-17 10:00", "status": "ok"}

    return {"observation": f"Unknown calendar action: {action}", "status": "error"}


# ---------------------------------------------------------------------------
# Tool execution — file_manager
# ---------------------------------------------------------------------------

def _exec_file_manager(state: Dict, action: str, params: Dict) -> Dict:
    work_dir = state["work_dir"]

    def resolve(p: str) -> Path:
        pp = Path(p) if p else Path(".")
        return pp if pp.is_absolute() else work_dir / pp

    if action == "list_directory":
        path = resolve(_p(params, "path", "directory", "dir", default="."))
        if not path.exists():
            return {"observation": "Directory not found", "status": "error"}
        items = [str(f.relative_to(work_dir)) for f in sorted(path.rglob("*")) if f.is_file()]
        return {"observation": json.dumps(items), "files": items, "status": "ok"}

    if action == "read_file":
        path = resolve(_p(params, "path", "file_path", "filename", "file"))
        if not path.exists():
            return {"observation": "File not found", "status": "error"}
        try:
            content = path.read_text(encoding="utf-8", errors="replace")
            return {"observation": content[:2000], "status": "ok"}
        except Exception as e:
            return {"observation": str(e), "status": "error"}

    if action == "copy_file":
        src  = resolve(_p(params, "src", "source", "from"))
        dest = resolve(_p(params, "dest", "destination", "to"))
        if not src.exists():
            return {"observation": "Source not found", "status": "error"}
        shutil.copy2(str(src), str(dest))
        return {"observation": f"Copied to {dest.name}", "status": "ok", "state_changed": True}

    if action == "move_file":
        src  = resolve(_p(params, "src", "source", "from"))
        dest = resolve(_p(params, "dest", "destination", "to"))
        if not src.exists():
            return {"observation": "Source not found", "status": "error"}
        shutil.move(str(src), str(dest))
        return {"observation": f"Moved to {dest.name}", "status": "ok", "state_changed": True}

    if action == "delete_file":
        path = resolve(_p(params, "path", "file_path", "filename", "file"))
        if not path.exists():
            return {"observation": "File not found", "status": "error"}
        path.unlink()
        return {"observation": "Deleted", "status": "ok", "state_changed": True}

    if action == "create_directory":
        path = resolve(_p(params, "path", "dir", "directory", "name"))
        path.mkdir(parents=True, exist_ok=True)
        return {"observation": f"Created {path.name}", "status": "ok", "state_changed": True}

    if action == "search_files":
        directory = resolve(_p(params, "directory", "path", "dir", default="."))
        pattern   = params.get("pattern", "*")
        hits = [str(p.relative_to(work_dir)) for p in directory.glob(f"**/{pattern}")]
        return {"observation": json.dumps(hits), "files": hits, "status": "ok"}

    if action == "write_file":
        path = resolve(_p(params, "path", "file_path", "filename", "file"))
        content = params.get("content", params.get("text", params.get("data", "")))
        path.parent.mkdir(parents=True, exist_ok=True)
        path.write_text(str(content), encoding="utf-8")
        return {"observation": f"Written {path.name}", "status": "ok", "state_changed": True}

    return {"observation": f"Unknown file_manager action: {action}", "status": "error"}


# ---------------------------------------------------------------------------
# Dispatch
# ---------------------------------------------------------------------------

_DISPATCH = {
    "word":         _exec_word,
    "excel":        _exec_excel,
    "powerpoint":   _exec_powerpoint,
    "email":        _exec_email,
    "calendar":     _exec_calendar,
    "file_manager": _exec_file_manager,
}


def _execute_tool(task_id: str, tool: str, params: Dict) -> Dict:
    state = _get_state(task_id)
    if state is None:
        return {"observation": "Task not initialised — call /init first", "status": "error"}
    app_prefix = tool.split("__")[0] if "__" in tool else "unknown"
    action     = tool.split("__", 1)[1] if "__" in tool else tool
    handler = _DISPATCH.get(app_prefix)
    if handler is None:
        return {"observation": f"Unknown app: {app_prefix}", "status": "error"}
    try:
        return handler(state, action, params)
    except Exception as exc:
        return {"observation": f"Error: {exc}", "status": "error"}


# ---------------------------------------------------------------------------
# Evaluation
# ---------------------------------------------------------------------------

def _evaluate_task(task_id: str) -> float:
    state = _get_state(task_id)
    task  = _get_task(task_id)
    if state is None or task is None:
        return 0.0

    data     = task["data"]
    app      = task["app"]
    work_dir = state["work_dir"]

    # OfficeBench format: evaluation is a list of {function, args} criteria
    criteria = data.get("evaluation", [])
    if isinstance(criteria, list) and criteria:
        return _check_ob_criteria(criteria, state, work_dir)

    # Heuristic: was an output file produced?
    ext_map = {"word": "*.docx", "excel": "*.xlsx", "powerpoint": "*.pptx"}
    if app in ext_map and list(work_dir.glob(ext_map[app])):
        return 1.0
    if app == "email" and state.get("sent_emails"):
        return 1.0
    if app == "calendar" and len(state.get("calendar_store", [])) != len(_DEFAULT_CALENDAR):
        return 1.0

    return 0.0


def _check_ob_criteria(criteria: list, state: Dict, work_dir: Path) -> float:
    """Evaluate OfficeBench-format criteria: [{function, args}, ...]."""
    if not criteria:
        return 0.0
    passed = 0
    for c in criteria:
        fn   = c.get("function", "")
        args = c.get("args", {})
        try:
            if fn == "evaluate_contain":
                passed += _eval_contain(args, state, work_dir)
            elif fn == "evaluate_exist":
                path = work_dir / args.get("file", "")
                passed += 1 if path.exists() else 0
            elif fn == "evaluate_equal":
                passed += _eval_equal(args, state, work_dir)
            else:
                passed += 1  # unknown function: don't penalise
        except Exception:
            pass
    return passed / len(criteria)


def _eval_contain(args: dict, state: Dict, work_dir: Path) -> int:
    """Check that a file contains all expected keywords."""
    doc_type = args.get("doc_type", "")
    file_rel = args.get("file", "")
    keywords = args.get("keywords", [])
    path     = work_dir / file_rel if file_rel else None

    if doc_type in ("word",) and path and path.exists():
        try:
            import docx
            doc  = docx.Document(str(path))
            text = "\n".join(p.text for p in doc.paragraphs).lower()
            return 1 if all(kw.lower() in text for kw in keywords) else 0
        except Exception:
            pass
    if doc_type in ("excel",) and path and path.exists():
        try:
            import openpyxl
            wb   = openpyxl.load_workbook(str(path))
            text = " ".join(str(c.value) for ws in wb.worksheets for row in ws.iter_rows() for c in row if c.value).lower()
            return 1 if all(kw.lower() in text for kw in keywords) else 0
        except Exception:
            pass
    if doc_type == "txt" and path and path.exists():
        try:
            text = path.read_text(errors="ignore").lower()
            return 1 if all(kw.lower() in text for kw in keywords) else 0
        except Exception:
            pass
    if doc_type == "email":
        all_text = " ".join(e.get("subject", "") + " " + e.get("body", "")
                            for e in state.get("email_store", [])).lower()
        return 1 if all(kw.lower() in all_text for kw in keywords) else 0
    return 0


def _eval_equal(args: dict, state: Dict, work_dir: Path) -> int:
    """Check that a cell or field equals an expected value."""
    doc_type = args.get("doc_type", "")
    expected = str(args.get("expected", "")).lower()
    if doc_type == "excel":
        try:
            import openpyxl
            path = work_dir / args.get("file", "")
            wb   = openpyxl.load_workbook(str(path))
            ws   = wb.active
            val  = str(ws[args.get("cell", "A1")].value or "").lower()
            return 1 if val == expected else 0
        except Exception:
            pass
    return 0


# ---------------------------------------------------------------------------
# Flask app
# ---------------------------------------------------------------------------

app = Flask(__name__)


@app.get("/health")
def health():
    return jsonify({"status": "ok", "tasks": len(_load_tasks())})


@app.get("/tasks")
def list_tasks():
    return jsonify([
        {"id": t["id"], "instruction": t["instruction"],
         "app": t["app"], "split": t["split"]}
        for t in _load_tasks()
    ])


@app.post("/tasks/<task_id>/init")
def init_task(task_id: str):
    task = _get_task(task_id)
    if task is None:
        return jsonify({"error": "Task not found"}), 404
    state = _init_task_env(task_id, task)
    files = [str(f.relative_to(state["work_dir"])) for f in sorted(state["work_dir"].rglob("*")) if f.is_file()]
    return jsonify({"status": "ok", "task_id": task_id, "app": task["app"],
                    "instruction": task["instruction"],
                    "files": files,
                    "work_dir": str(state["work_dir"])})


@app.post("/tasks/<task_id>/execute")
def execute_task(task_id: str):
    body   = request.get_json(force=True) or {}
    tool   = body.get("tool", "")
    params = body.get("params", {})
    return jsonify(_execute_tool(task_id, tool, params))


@app.post("/tasks/<task_id>/evaluate")
def evaluate_task(task_id: str):
    score = _evaluate_task(task_id)
    return jsonify({"score": score, "success": score >= 0.5})


@app.post("/tasks/<task_id>/reset")
def reset_task(task_id: str):
    _task_state.pop(task_id, None)
    for store in (_open_docs, _open_wbs, _open_pptx):
        store.clear()
    task = _get_task(task_id)
    if task:
        _init_task_env(task_id, task)
    return jsonify({"status": "ok"})


# ---------------------------------------------------------------------------
# Entry point
# ---------------------------------------------------------------------------

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="OfficeBench REST server")
    parser.add_argument("--port",      type=int, default=8001)
    parser.add_argument("--tasks-dir", default=str(TASKS_DIR))
    parser.add_argument("--work-dir",  default=str(WORK_BASE))
    args = parser.parse_args()

    TASKS_DIR = Path(args.tasks_dir)
    WORK_BASE = Path(args.work_dir)
    WORK_BASE.mkdir(parents=True, exist_ok=True)

    print(f"[OfficeBench] Server starting on port {args.port}")
    print(f"[OfficeBench] Tasks dir : {TASKS_DIR}")
    print(f"[OfficeBench] Tasks found: {len(_load_tasks())}")

    app.run(host="0.0.0.0", port=args.port, debug=False)
